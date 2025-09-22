# Copyright (c) Meta Platforms, Inc. and affiliates.
# All rights reserved.
#
# This source code is licensed under the terms described in the LICENSE file in
# the root directory of this source tree.

# Based on https://github.com/microsoft/markitdown/blob/main/packages/markitdown/src/markitdown/converters/_markdownify.py
# Licence: https://github.com/microsoft/markitdown/blob/main/LICENSE

import base64
import binascii
import copy
import html
import io
import mimetypes
import os
import re
import tempfile
import traceback
import xml.etree.ElementTree as ET
from itertools import combinations
from typing import IO, Any
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx
import puremagic
import requests
from bs4 import BeautifulSoup
from requests.models import Response

from are.simulation.exceptions import MarkdownConverterError

INTERACTIONS_BY_TAG = {
    "button": "click",
    "details": "click",
    "label": "click",
    "menu": "click",
    "menuitem": "click",
    "select": "select",
    "textarea": "input_text",
    "summary": "click",
}

INTERACTIVE_ROLES = {
    "button",
    "menu",
    "menuitem",
    "link",
    "checkbox",
    "radio",
    "slider",
    "tab",
    "tabpanel",
    "textbox",
    "combobox",
    "grid",
    "listbox",
    "option",
    "progressbar",
    "scrollbar",
    "searchbox",
    "switch",
    "tree",
    "treeitem",
    "spinbutton",
    "tooltip",
}

INTERACTIVE_ARIA_ROLES = {
    "menu",
    "menuitem",
    "button",
}

HIDDEN_CLASSES = {
    "hidden",
    "invisible",
    "d-none",
    "sr-only",
    "collapse",
    "fade",
    "modal-hidden",
}

FILTER_ATTRIBUTES = {
    "id",
    "name",
    "placeholder",
    "title",
    "href",
    "src",
    "aria-label",
    "aria-name",
    "aria-role",
    "aria-description",
    "aria-expanded",
    "aria-haspopup",
    "type",
    "value",
    "required",
}


class FileConversionMarkdownError(MarkdownConverterError):
    """
    The file format is accepted by a Markdown converter, but an error occurred during conversion.
    """

    pass


class UnsupportedFormatMarkdownError(MarkdownConverterError):
    """
    The file format could not be rendered by any Markdown converter.
    """

    pass


class CustomMarkdownConverter(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    - Add interactive elements in simplified HTML format to allow for interaction (if self.interactive)
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        self.interactive = options.get("interactive", False)
        self.suggest_actions = options.get("suggest_actions", True)
        self.interactive_elements = []
        self.element_counter = 0
        super().__init__(**options)

    def convert_soup(self, soup: Any) -> str:
        self.element_counter = 0  # reset on every page update
        self.interactive_elements = []
        markdown = super().convert_soup(soup)  # type: ignore
        for elem_info in self.interactive_elements:
            elem_info["css_selector"] = self.build_css_selector(elem_info)
        return markdown

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        try:
            prefix, suffix, text = markdownify.chomp(text)  # type: ignore
            if not text:
                return ""
            href = el.get("href")
            title = el.get("title")

            # Escape URIs and skip non-http or file schemes
            if href:
                try:
                    parsed_url = urlparse(href)  # type: ignore
                    if parsed_url.scheme and parsed_url.scheme.lower() not in [
                        "http",
                        "https",
                        "file",
                    ]:  # type: ignore
                        return "%s%s%s" % (prefix, text, suffix)
                    href = urlunparse(
                        parsed_url._replace(path=quote(unquote(parsed_url.path)))
                    )  # type: ignore
                except ValueError:  # It's not clear if this ever gets thrown
                    return "%s%s%s" % (prefix, text, suffix)

            # For the replacement see #29: text nodes underscores are escaped
            if (
                self.options["autolinks"]
                and text.replace(r"\_", "_") == href
                and not title
                and not self.options["default_title"]
            ):
                # Shortcut syntax
                return "<%s>" % href
            if self.options["default_title"] and not title:
                title = href
            title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
            return (
                "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
                if href
                else text
            )
        except Exception:
            return super().convert_a(el, text, convert_as_inline)  # type: ignore

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        # if (
        #    convert_as_inline
        #    and el.parent.name not in self.options["keep_inline_images_in"]
        # ):
        #    return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return f"![{alt}]({src}{title_part})"

    def convert_button(self, el: Any, text: str, convert_as_inline: bool):
        return self.convert_interactive_element("button", el, text)

    def convert_input(self, el: Any, text: str, convert_as_inline: bool) -> str:
        return self.convert_interactive_element("input", el, text)

    def convert_textarea(self, el: Any, text: str, convert_as_inline: bool) -> str:
        return self.convert_interactive_element("textarea", el, text)

    def convert_menu(self, el: Any, text: str, convert_as_inline: bool) -> str:
        return self.convert_interactive_element("menu", el, text)

    def convert_menuitem(self, el: Any, text: str, convert_as_inline: bool) -> str:
        return self.convert_interactive_element("menuitem", el, text)

    def convert_summary(self, el: Any, text: str, convert_as_inline: bool):
        return self.convert_interactive_element("summary", el, text)

    def convert_details(self, el: Any, text: str, convert_as_inline: bool):
        return self.convert_interactive_element("details", el, text)

    def convert_select(self, el: Any, text: str, convert_as_inline: bool) -> str:
        tag_name = "select"
        if (
            self.interactive
            and self._is_visible(el)
            and self._is_interactive(el, tag_name)
        ):
            try:
                attributes = self._get_attributes(tag_name, el, text)
                options = []
                for option in el.find_all("option"):
                    option_text = option.get_text(strip=True)
                    option_attributes = self._get_attributes(
                        "option", option, option_text
                    )
                    options.append(option_attributes)
                attributes["options"] = options
                placeholder = self._get_placeholder(el, tag_name, attributes)
                if placeholder:
                    return placeholder
            except Exception:
                pass
        return text

    def _is_unique(self, candidate_locator: dict) -> bool:
        """Check attribute subset uniqueness"""
        matches = 0
        for other in self.interactive_elements:
            if all(
                other["attributes"].get(k) == v for k, v in candidate_locator.items()
            ):
                matches += 1
        return matches == 1

    def _is_visible(self, el: Any) -> bool:
        """Check if an element is visible to users."""
        style = el.get("style", "").lower()
        if any(
            x in style for x in ["display: none", "visibility: hidden", "opacity: 0"]
        ):
            return False
        classes = set(el.get("class", []))
        if classes & HIDDEN_CLASSES:
            return False
        if el.get("hidden") is not None:
            return False
        if el.get("aria-hidden", "").lower() == "true":
            return False
        return True

    def _is_interactive(self, el: Any, tag_name: str) -> bool:
        role = el.get("role", "").lower()
        return (
            tag_name in INTERACTIONS_BY_TAG
            or role in INTERACTIVE_ROLES
            or role in INTERACTIVE_ARIA_ROLES
            or el.get("onclick")
            or el.get("onchange")
            or el.get("contenteditable") == "true"
        )

    def convert_interactive_element(self, tag_name: str, el: Any, text: str):
        """Convert an interactive element to a simplified ArchiveSearchTool format."""
        if (
            self.interactive
            and self._is_visible(el)
            and self._is_interactive(el, tag_name)
        ):
            try:
                attributes = self._get_attributes(tag_name, el, text)
                placeholder = self._get_placeholder(el, tag_name, attributes)
                if placeholder:
                    return placeholder
            except Exception:
                pass
        return text

    def _suggest_action(self, tag_name: str, attributes: dict[str, Any]):
        """Determine the appropriate default action for an element."""
        # If element is contenteditable (explicitly "true" or empty string), assume text input.
        contenteditable = attributes.get("contenteditable", "").lower()
        if contenteditable == "true" or (
            contenteditable == "" and "contenteditable" in attributes
        ):
            return "input_text"

        # Textareas are clearly for text input.
        if tag_name == "textarea":
            return "input_text"

        # For input elements, decide based on the type.
        if tag_name == "input":
            input_type = attributes.get("type", "text").lower()
            # For checkbox or radio, use "check"
            if input_type in ["checkbox", "radio"]:
                return "check"
            # For submit or button inputs, clicking is appropriate.
            elif input_type in ["submit", "button"]:
                return "click"
            # For file inputs, allow file uploads.
            elif input_type == "file":
                return "upload_files"
            # Otherwise, assume a text input field.
            else:
                return "input_text"

        # For select elements, use select action.
        if tag_name == "select":
            return "select"

        # If the element has an explicit role attribute, adjust accordingly.
        role = attributes.get("role", "").lower()
        if role == "button":
            return "click"
        if role in ["textbox", "searchbox", "combobox"]:
            return "input_text"

        # Default fallback action is "click"
        return "click"

    def _get_attributes(self, tag_name: str, el: Any, text: str) -> dict[str, Any]:
        """Extract attributes to display from an interactive element."""
        attributes = {}
        for k in el.attrs:
            if k in FILTER_ATTRIBUTES:  # Filter relevant attributes
                attributes[k] = el.get(k)
            elif k.startswith("data-"):  # Collect data- attributes
                attributes[k] = el.get(k)
        # Don’t store text for container-like tags
        if tag_name not in ["details", "summary"]:
            attributes["text"] = text.replace("\n", " ").strip()
        return attributes

    def _get_placeholder(
        self, el: Any, tag_name: str, attributes: dict[str, str]
    ) -> str:
        """Show an interactive element in simplified HTML format."""
        idx = self.element_counter
        self.element_counter += 1
        if self.suggest_actions:
            action = self._suggest_action(tag_name, attributes)
            if action:
                attributes["allowed"] = action
        xpath = self._generate_xpath(el)
        attr_str = " ".join(
            "{}='{}'".format(k, re.sub(r"\s+", " ", str(v)).strip())
            for k, v in attributes.items()
            if v
        )
        placeholder = f"<{tag_name} idx={idx} {attr_str}/>"
        element_info = {
            "idx": idx,
            "tag_name": tag_name,
            "element": el,
            "attributes": attributes,
            "xpath": xpath,
            "placeholder": placeholder,
        }
        self.interactive_elements.append(element_info)
        return placeholder

    def _generate_xpath(self, element):
        """Generate a unique XPath for an element."""
        components = []
        current = element
        for parent in current.parents:
            siblings = parent.find_all(current.name, recursive=False)
            if len(siblings) > 1:
                index = siblings.index(current) + 1
                components.append(f"{current.name}[{index}]")
            else:
                components.append(current.name)
            current = parent
        components.reverse()
        xpath = "//" + "/".join(components)
        return xpath

    def _select_attributes(self, elem_info) -> dict:
        """
        Return a dict of (attribute -> value) that uniquely identifies the element
        among all interactive elements, giving priority to more stable attributes.
        """
        raw_attrs = elem_info["attributes"]
        # Data attributes are volatile; do not use them as locators
        candidate_keys = [
            k for k, v in raw_attrs.items() if not k.startswith("data-") and v
        ]

        # Define a prioritized list of stable attributes
        stable_priority = [
            "id",
            "name",
            "aria-label",
            "aria-labelledby",
            "placeholder",
            "title",
        ]

        # Sort candidate_keys so that keys appearing in stable_priority come first
        candidate_keys.sort(
            key=lambda k: (
                stable_priority.index(k)
                if k in stable_priority
                else len(stable_priority)
            )
        )

        # Try each candidate key on its own first.
        for key in candidate_keys:
            candidate_locator = {key: raw_attrs[key]}
            if self._is_unique(candidate_locator):
                return candidate_locator

        # If no single attribute is unique, try combinations (starting from size 2)
        for size in range(2, len(candidate_keys) + 1):
            for subset in combinations(candidate_keys, size):
                candidate_locator = {k: raw_attrs[k] for k in subset}
                if self._is_unique(candidate_locator):
                    return candidate_locator

        return {}

    def build_css_selector(self, elem_info: dict) -> str:
        """
        Create a unique Playwright element locator using a stable set of attributes.
        Priority:
        1. If an 'id' attribute exists, use that.
        2. Otherwise, use a minimal set of stable attributes.
        3. Finally, fall back to a generated XPath.
        """
        attributes = elem_info["attributes"]
        # (1) Use the `id` attribute if present
        elem_id = attributes.get("id")
        if elem_id:
            return f"#{elem_id}"

        # 2. Build selector with tag and minimal unique attributes
        candidate_attrs = {
            k: v
            for k, v in attributes.items()
            if k in FILTER_ATTRIBUTES and v and not k.startswith("data-")
        }

        # Prioritize stable attributes
        priority_attrs = ["id", "name", "aria-label", "type", "placeholder"]
        tag_name = elem_info["tag_name"]
        selector_parts = [tag_name]

        for attr in priority_attrs:
            if attr in candidate_attrs:
                value = candidate_attrs[attr].replace('"', '\\"')
                selector_parts.append(f'[{attr}="{value}"]')
                locator = "".join(selector_parts)
                if self._is_unique({attr: candidate_attrs[attr]}):
                    return locator

        # 3. Add more attributes if needed
        for attr, value in candidate_attrs.items():
            if attr not in priority_attrs:
                value = value.replace('"', '\\"')
                selector_parts.append(f'[{attr}="{value}"]')
                locator = "".join(selector_parts)
                if self._is_unique(
                    {
                        k: candidate_attrs[k]
                        for k in candidate_attrs.keys()
                        if f'[{k}="' in locator
                    }
                ):
                    return locator

        # 4. Fallback to text content if unique
        text = attributes.get("text")
        if text and self._is_unique({"text": text}):
            clean_text = text.replace('"', '\\"')
            return f'{tag_name}:has-text("{clean_text}")'

        # (3) Fallback to positional XPath (this is brittle)
        xpath = self._generate_xpath(elem_info["element"])
        return f"xpath={xpath}"


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(
        self,
        title: str | None = None,
        text_content: str = "",
        interactive_elements=None,
    ):
        self.title = title
        self.text_content = text_content
        self.interactive_elements = interactive_elements or []


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert_io(
        self, file_io: IO[bytes], **kwargs: Any
    ) -> DocumentConverterResult | None:
        """Convert from an IO object (primary method)"""
        raise NotImplementedError()

    def convert_path(
        self, local_path: str, **kwargs: Any
    ) -> DocumentConverterResult | None:
        """Convert from local path (backward compatibility)"""
        with open(local_path, "rb") as f:
            return self.convert_io(f, **kwargs)

    # Deprecated - kept for backward compatibility
    def convert(self, local_path: str, **kwargs: Any) -> DocumentConverterResult | None:
        return self.convert_path(local_path, **kwargs)


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type("__placeholder" + (file_extension or ""))

        # Only accept text files
        if content_type is None:
            return None
        # elif "text/" not in content_type.lower():
        #    return None

        text_content = ""
        # Read from IO object and decode as UTF-8
        file_io.seek(0)  # Ensure we're at the beginning
        content_bytes = file_io.read()
        text_content = content_bytes.decode("utf-8")

        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def __init__(self, interactive: bool = False):
        self.interactive = interactive
        self.markdownify = CustomMarkdownConverter(interactive=self.interactive)

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        extension = file_extension or ""
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        file_io.seek(0)  # Ensure we're at the beginning
        content_bytes = file_io.read()

        for encoding in ["utf-8", "ISO-8859-1"]:
            try:
                html_content = content_bytes.decode(encoding)
                result = self._convert(html_content)
                break
            except Exception:
                pass
        return result

    def _convert(self, html_content) -> DocumentConverterResult | None:
        """Helper function that converts and HTML string."""
        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = self.markdownify.convert_soup(body_elm)
        else:
            webpage_text = self.markdownify.convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
            interactive_elements=self.markdownify.interactive_elements,
        )


class BingSerpConverter(DocumentConverter):
    """
    Handle Bing results pages (only the organic search results).
    NOTE: It is better to use the Bing API
    """

    def convert(self, local_path, **kwargs) -> DocumentConverterResult | None:
        # Bail if not a Bing SERP
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        url = kwargs.get("url", "")
        if not re.search(r"^https://www\.bing\.com/search\?q=", url):
            return None

        # Parse the query parameters
        parsed_params = parse_qs(urlparse(url).query)
        query = parsed_params.get("q", [""])[0]

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Clean up some formatting
        for tptt in soup.find_all(class_="tptt"):
            if hasattr(tptt, "string") and tptt.string:  # type: ignore
                tptt.string += " "  # type: ignore
        for slug in soup.find_all(class_="algoSlug_icon"):
            slug.extract()

        # Parse the algorithmic results
        _markdownify = CustomMarkdownConverter()
        results = list()
        for result in soup.find_all(class_="b_algo"):
            # Rewrite redirect urls
            for a in result.find_all("a", href=True):  # type: ignore
                parsed_href = urlparse(a["href"])  # type: ignore
                qs = parse_qs(parsed_href.query)

                # The destination is contained in the u parameter,
                # but appears to be base64 encoded, with some prefix
                if "u" in qs:
                    u = (
                        qs["u"][0][2:].strip() + "=="
                    )  # Python 3 doesn't care about extra padding

                    try:
                        # RFC 4648 / Base64URL" variant, which uses "-" and "_"
                        a["href"] = base64.b64decode(u, altchars="-_").decode("utf-8")  # type: ignore
                    except UnicodeDecodeError:
                        pass
                    except binascii.Error:
                        pass

            # Convert to markdown
            md_result = _markdownify.convert_soup(result).strip()
            lines = [line.strip() for line in re.split(r"\n+", md_result)]
            results.append("\n".join([line for line in lines if len(line) > 0]))

        webpage_text = (
            f"## A Bing search for '{query}' found the following results:\n\n"
            + "\n\n".join(results)
        )

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Bail if not a PDF
        extension = file_extension or ""
        if extension.lower() != ".pdf":
            return None

        # pdfminer doesn't support IO[bytes] directly, so we need to use a temporary file
        # This is a fallback for libraries that don't support IO objects
        file_io.seek(0)  # Ensure we're at the beginning
        content = file_io.read()

        # Use context manager for temporary file
        with tempfile.NamedTemporaryFile(suffix=".pdf") as temp_file:
            temp_file.write(content)
            temp_file.flush()  # Ensure content is written to disk
            text_content = pdfminer.high_level.extract_text(temp_file.name)
            return DocumentConverterResult(
                title=None,
                text_content=text_content,
            )


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Bail if not a DOCX
        extension = file_extension or ""
        if extension.lower() != ".docx":
            return None

        file_io.seek(0)  # Ensure we're at the beginning
        result = mammoth.convert_to_html(file_io)
        html_content = result.value
        return self._convert(html_content)


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Bail if not a XLSX
        extension = file_extension or ""
        if extension.lower() not in [
            ".xlsx",
            ".xls",
            ".xlsm",
            ".xlsb",
            ".ods",
            ".odf",
            ".odt",
        ]:
            return None

        file_io.seek(0)  # Ensure we're at the beginning
        sheets = pd.read_excel(file_io, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            converted_result = self._convert(html_content)
            if converted_result is not None:
                md_content += converted_result.text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class XmlConverter(DocumentConverter):
    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Bail if not a XML
        extension = file_extension or ""
        if extension.lower() != ".xml":
            return None

        file_io.seek(0)  # Ensure we're at the beginning
        content_bytes = file_io.read()
        xml_string = content_bytes.decode("utf-8")

        return self._convert_xml_string(xml_string)

    def _convert_xml_string(self, xml_string: str) -> DocumentConverterResult | None:
        def extract_table_from_html_like(xml_root):
            table = xml_root.find(".//table")
            if table is None:
                raise ValueError("No table found in the XML")

            headers = [th.text for th in table.find("thead").findall("th")]
            rows = [
                [td.text for td in tr.findall("td")]
                for tr in table.find("tbody").findall("tr")
            ]

            # Create markdown table
            markdown = "| " + " | ".join(headers) + " |\n"
            markdown += "| " + " | ".join(["---"] * len(headers)) + " |\n"
            for row in rows:
                markdown += "| " + " | ".join(row) + " |\n"

        def extract_table_from_wordml(xml_root, namespaces):
            # Parse the XML content
            root = xml_root
            namespace = {"w": "http://schemas.microsoft.com/office/word/2003/wordml"}

            # Extract text content
            body = root.find("w:body", namespace)
            paragraphs = body.findall(".//w:p", namespace)
            text_content = []
            for para in paragraphs:
                texts = para.findall(".//w:t", namespace)
                for text in texts:
                    text_content.append(text.text)

            return "\n".join(text_content)

        # Parse the XML string
        root = ET.fromstring(xml_string)
        namespaces = {"w": "http://schemas.microsoft.com/office/word/2003/wordml"}

        if root.tag.endswith("wordDocument"):
            markdown = extract_table_from_wordml(root, namespaces)
        else:
            markdown = extract_table_from_html_like(root)

        return DocumentConverterResult(
            title=None,
            text_content=markdown.strip() if markdown else "",
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert_io(
        self, file_io: IO[bytes], file_extension: str | None = None, **kwargs: Any
    ) -> DocumentConverterResult | None:
        # Bail if not a PPTX
        extension = file_extension or ""
        if extension.lower() != ".pptx":
            return None

        file_io.seek(0)  # Ensure we're at the beginning
        presentation = pptx.Presentation(file_io)
        return self._convert_presentation(presentation)

    def _convert_presentation(self, presentation) -> DocumentConverterResult | None:
        md_content = ""
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    # try:
                    #    filename = shape.image.filename
                    # except:
                    #    pass

                    md_content += (
                        "\n!["
                        + (alt_text if alt_text else shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:  # type: ignore
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    converted_result = self._convert(html_table)
                    if converted_result is not None:
                        md_content += (
                            "\n" + converted_result.text_content.strip() + "\n"
                        )

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"  # type: ignore
                    else:
                        md_content += shape.text + "\n"  # type: ignore

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:  # type: ignore
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:  # type: ignore
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:  # type: ignore
            return True
        return False


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        interactive: bool = False,
    ):
        self.interactive = interactive

        self._page_converters: list[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.add_converter(PdfConverter())
        self.add_converter(XmlConverter())
        self.add_converter(PptxConverter())
        self.add_converter(XlsxConverter())
        self.add_converter(DocxConverter())
        self.add_converter(BingSerpConverter())
        self.add_converter(HtmlConverter(interactive=self.interactive))
        self.add_converter(PlainTextConverter())

    def add_converter(self, converter: DocumentConverter):
        self._page_converters.append(converter)

    def convert_io(self, file_io: IO[bytes], **kwargs: Any) -> DocumentConverterResult:
        """Convert from IO object"""
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Convert
        return self._convert_io(file_io, extensions, **kwargs)

    def convert_path(self, local_path: str, **kwargs: Any) -> DocumentConverterResult:
        """Convert from local file path"""
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(local_path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(local_path))

        # Convert
        return self._convert(local_path, extensions, **kwargs)

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:
        """Convert from URL"""
        # Send a HTTP request to the URL
        response = requests.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: Response, **kwargs: Any
    ) -> DocumentConverterResult:
        """Convert from HTTP response"""
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Download content and convert directly from IO
        content = b""
        for chunk in response.iter_content(chunk_size=512):
            content += chunk

        content_io = io.BytesIO(content)
        return self._convert_io(content_io, extensions, url=response.url, **kwargs)

    def convert_str(self, content: str, **kwargs: Any) -> DocumentConverterResult:
        """Convert from string content"""
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Convert string to bytes
        if isinstance(content, str):
            content_bytes = content.encode("utf-8")
        else:
            content_bytes = content

        # Convert directly from IO
        content_io = io.BytesIO(content_bytes)
        return self._convert_io(content_io, extensions, **kwargs)

    def convert_local(self, local_path: str, **kwargs: Any) -> DocumentConverterResult:
        """Convert from local file path (deprecated, use convert_path instead)"""
        return self.convert_path(local_path, **kwargs)

    # Deprecated - kept for backward compatibility
    def convert(
        self, source: str | Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a Response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_path(source, **kwargs)
        # Request response
        elif isinstance(source, Response):
            return self.convert_response(source, **kwargs)

    def _convert_io(
        self, file_io: IO[bytes], extensions: list[str | None], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                try:
                    res = converter.convert_io(file_io, **_kwargs)
                    if res is not None:
                        # Normalize the content
                        res.text_content = "\n".join(
                            [
                                line.rstrip()
                                for line in re.split(r"\r?\n", res.text_content)
                            ]
                        )
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                        return res

                # Catch any uncaught exceptions
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionMarkdownError(
                f"Could not convert file to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatMarkdownError(
            f"Could not convert file to Markdown. The formats {extensions} are not supported."
        )

    def _convert(
        self, local_path: str, extensions: list[str | None], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                try:
                    res = converter.convert(local_path, **_kwargs)
                    if res is not None:
                        # Normalize the content
                        res.text_content = "\n".join(
                            [
                                line.rstrip()
                                for line in re.split(r"\r?\n", res.text_content)
                            ]
                        )
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                        return res

                # Catch any uncaught exceptions
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionMarkdownError(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatMarkdownError(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None
